import sys
sys.path.append('/Users/madelynzhu/Library/Python/3.13/lib/python/site-packages')
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_BG = RGBColor(9, 13, 22)         # #090D16 Dark Navy
    COLOR_CARD = RGBColor(19, 27, 46)      # #131B2E Dark Card
    COLOR_TEXT = RGBColor(248, 250, 252)   # #F8FAFC White Text
    COLOR_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Muted Text
    COLOR_CYAN = RGBColor(56, 189, 248)    # #38BDF8 Accent Cyan
    COLOR_EMERALD = RGBColor(52, 211, 153) # #34D399 Accent Emerald
    COLOR_AMBER = RGBColor(251, 191, 36)   # #FBBF24 Accent Amber
    COLOR_ROSE = RGBColor(244, 63, 94)     # #F43F5E Accent Rose

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, category, title, subtitle):
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN
        p_cat.font.name = "Arial"

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT
        p_title.font.name = "Arial"

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = COLOR_MUTED
        p_sub.font.name = "Arial"

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Title Card Box
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(3.0))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "🔬 Biosensor Model Training Initiative"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_CYAN

    p2 = tf1.add_paragraph()
    p2.text = "100-Strain Bacterial Data Catalog"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "Automated Data Extraction, Entity Resolution, Code Evolution, and Biosafety Compliance"
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(10)

    # 3 Metadata Cards (Bottom)
    card_data = [
        ("Code Evolution", "4 Major GitHub Commit Iterations", COLOR_CYAN),
        ("Project Scope", "100 Target Strains Cataloged", COLOR_EMERALD),
        ("Primary Deliverable", "37-Column Populated Schema", COLOR_AMBER)
    ]
    for idx, (label, val, accent) in enumerate(card_data):
        left_pos = Inches(1.2 + idx * 3.7)
        card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(4.8), Inches(3.4), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = accent
        
        ctf = card.text_frame
        ctf.word_wrap = True
        cp1 = ctf.paragraphs[0]
        cp1.text = label.upper()
        cp1.font.size = Pt(10)
        cp1.font.bold = True
        cp1.font.color.rgb = COLOR_MUTED
        
        cp2 = ctf.add_paragraph()
        cp2.text = val
        cp2.font.size = Pt(14)
        cp2.font.bold = True
        cp2.font.color.rgb = COLOR_TEXT
        cp2.space_before = Pt(8)

    # Speaker Notes Slide 1
    slide1.notes_slide.notes_text_frame.text = (
        "Hello everyone. Today I'm presenting the project report for our 100-strain bacterial reference catalog. "
        "As the data analyst on this project, I developed the codebase and data mining pipeline from initial prototype "
        "to a production 37-column dataset for training our disease biosensor. In this presentation, I'll walk through "
        "our code development history on GitHub, our technical data engineering architecture, and key microbiological findings."
    )

    # ==========================================
    # SLIDE 2: GitHub Code Evolution
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Slide 2 / 5 • Code Evolution & GitHub Updates", "Code Development Milestones & GitHub History", "Key engineering iterations, codebase refactoring, and feature enhancements across the repository lifecycle.")

    # Left Column: Commit History Card
    left_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(6.2), Inches(4.7))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR_CARD
    left_card.line.color.rgb = COLOR_CYAN

    ltf = left_card.text_frame
    ltf.word_wrap = True
    lp1 = ltf.paragraphs[0]
    lp1.text = "🌿 GITHUB DEVELOPMENT HISTORY"
    lp1.font.size = Pt(14)
    lp1.font.bold = True
    lp1.font.color.rgb = COLOR_CYAN

    commits = [
        ("v1.0 — Initial Scraper & Trial", "Built atcc_scraper.py using urllib & BS4. Conducted trial parsing on single-strain files."),
        ("v2.0 — 37-Column Script (atcc_scraper_Comp.py)", "Scaled extraction to 37 metadata attributes, added ThreadPoolExecutor & disk caching."),
        ("v2.5 — GenBank Link Transformation", "Converted raw accession text into direct NCBI URLs (ncbi.nlm.nih.gov/nuccore/)."),
        ("v3.0 — Scaling & Dataset Expansion", "Updated pipeline to process 148+ entries and refactored scripts for 200 Bug Project(FB).csv.")
    ]
    for c_title, c_desc in commits:
        cp_t = ltf.add_paragraph()
        cp_t.text = "• " + c_title
        cp_t.font.size = Pt(12)
        cp_t.font.bold = True
        cp_t.font.color.rgb = COLOR_TEXT
        cp_t.space_before = Pt(8)

        cp_d = ltf.add_paragraph()
        cp_d.text = "  " + c_desc
        cp_d.font.size = Pt(11)
        cp_d.font.color.rgb = COLOR_MUTED

    # Right Column: 4 Technical Upgrade Boxes
    upgrades = [
        ("NBF Fallback Engine", "Organism species lookup mapping (S. aureus ➔ ATCC 25923)", COLOR_EMERALD),
        ("GenBank API Links", "Automatic regex parsing of raw text into NCBI REST URLs", COLOR_CYAN),
        ("Cache Engine", "Local HTML response caching for sub-second re-runs", COLOR_AMBER),
        ("SSL Bypass", "macOS local CA certificate validation fix", COLOR_ROSE)
    ]
    for idx, (utitle, udesc, ucolor) in enumerate(upgrades):
        row = idx // 2
        col = idx % 2
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3 + col * 2.7), Inches(2.2 + row * 2.45), Inches(2.55), Inches(2.25))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = ucolor

        utf = card.text_frame
        utf.word_wrap = True
        up1 = utf.paragraphs[0]
        up1.text = utitle
        up1.font.size = Pt(13)
        up1.font.bold = True
        up1.font.color.rgb = ucolor

        up2 = utf.add_paragraph()
        up2.text = udesc
        up2.font.size = Pt(10)
        up2.font.color.rgb = COLOR_MUTED
        up2.space_before = Pt(6)

    # Speaker Notes Slide 2
    slide2.notes_slide.notes_text_frame.text = (
        "Here on Slide 2, I've highlighted our major code updates and GitHub commits: We started with version 1.0, "
        "building a baseline scraper in Python using BeautifulSoup and urllib, and performing trial extractions on "
        "single-strain files. In version 2.0, I upgraded to atcc_scraper_Comp.py, expanding output schema to 37 columns "
        "and adding ThreadPoolExecutor multithreading with local disk caching to prevent rate-limiting. In version 2.5, "
        "I converted raw GenBank text into direct clickable NCBI accession URLs. Finally, in version 3.0, we scaled "
        "the pipeline up to handle 148+ strain entries for the 200 Bug Project."
    )

    # ==========================================
    # SLIDE 3: Scraping Engine & Entity Resolution
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Slide 3 / 5 • CS & Technical Architecture", "Scraping Engine & Entity Resolution", "Algorithmic mapping of missing catalog identifiers, null value management, and building clean feature stores for ML models.")

    # 3 Architecture Cards
    arch_cards = [
        ("⚡ Parallel Architecture", "Multithreaded fetching with disk caching keeps latency under 0.5 seconds while preventing server rate limiting.", COLOR_CYAN),
        ("🔍 Taxonomic Entity Resolution", "Automated resolution of 'NBF' ('Not Being Found') species to primary ATCC reference catalog strains (e.g. S. aureus ➔ ATCC 25923).", COLOR_AMBER),
        ("🤖 37-Column Schema Output", "Generates a populated metadata catalog that acts as a structured feature store for biosensor machine learning models.", COLOR_EMERALD)
    ]
    for idx, (atitle, adesc, acolor) in enumerate(arch_cards):
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.95), Inches(2.2), Inches(3.7), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = acolor

        atf = card.text_frame
        atf.word_wrap = True
        ap1 = atf.paragraphs[0]
        ap1.text = atitle
        ap1.font.size = Pt(14)
        ap1.font.bold = True
        ap1.font.color.rgb = acolor

        ap2 = atf.add_paragraph()
        ap2.text = adesc
        ap2.font.size = Pt(11)
        ap2.font.color.rgb = COLOR_MUTED
        ap2.space_before = Pt(10)

    # Impact Callout Box
    impact_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.35))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = COLOR_CARD
    impact_box.line.color.rgb = COLOR_EMERALD

    itf = impact_box.text_frame
    itf.word_wrap = True
    ip1 = itf.paragraphs[0]
    ip1.text = "🧠 Technical Engineering Impact"
    ip1.font.size = Pt(12)
    ip1.font.bold = True
    ip1.font.color.rgb = COLOR_EMERALD

    ip2 = itf.add_paragraph()
    ip2.text = "Taxonomic entity resolution reduced catalog ambiguity from 22% in raw input to 0% in the final analytical dataset, ensuring clean join keys for machine learning model training."
    ip2.font.size = Pt(12)
    ip2.font.color.rgb = COLOR_TEXT
    ip2.space_before = Pt(4)

    # Speaker Notes Slide 3
    slide3.notes_slide.notes_text_frame.text = (
        "Zooming into technical feature highlights on Slide 3: my parallel web scraper uses disk caching to maintain "
        "sub-second re-run performance. A major engineering achievement was taxonomic entity resolution: automatically "
        "resolving missing 'NBF' entries to primary reference ATCC catalog numbers—such as mapping Staphylococcus aureus "
        "to ATCC 25923—which reduced dataset catalog ambiguity from 22% down to 0%."
    )

    # ==========================================
    # SLIDE 4: Microbiological Isolation Sources
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Slide 4 / 5 • Microbiological Domain", "Isolation Sources & Ecological Provenance", "Mapping isolation origins to establish sample diversity, target specificity, and real-world background matrix conditions for sensor testing.")

    # 3 Ecological Provenance Cards
    eco_cards = [
        ("🩸 Clinical & Human Isolates", "Isolated from blood, sputum, wound exudates, and tissue. Critical for benchmarking sensor sensitivity on human clinical pathogens (e.g., S. aureus, K. pneumoniae).", COLOR_ROSE),
        ("🐾 Veterinary Host Specimens", "Derived from animal host lesions and veterinary infections. Expands sensor target recognition across zoonotic vectors (e.g., Rhodococcus equi from foal lung abscess).", COLOR_AMBER),
        ("🌱 Environmental & Sewage", "Sampled from soil, wastewater, and sludge. Essential for testing sensor cross-reactivity against background flora (e.g., Corynebacterium glutamicum).", COLOR_EMERALD)
    ]
    for idx, (etitle, edesc, ecolor) in enumerate(eco_cards):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 3.95), Inches(2.2), Inches(3.7), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = ecolor

        etf = card.text_frame
        etf.word_wrap = True
        ep1 = etf.paragraphs[0]
        ep1.text = etitle
        ep1.font.size = Pt(14)
        ep1.font.bold = True
        ep1.font.color.rgb = ecolor

        ep2 = etf.add_paragraph()
        ep2.text = edesc
        ep2.font.size = Pt(11)
        ep2.font.color.rgb = COLOR_MUTED
        ep2.space_before = Pt(10)

    # Matrix Relevance Box
    matrix_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.35))
    matrix_box.fill.solid()
    matrix_box.fill.fore_color.rgb = COLOR_CARD
    matrix_box.line.color.rgb = COLOR_CYAN

    mtf = matrix_box.text_frame
    mtf.word_wrap = True
    mp1 = mtf.paragraphs[0]
    mp1.text = "💡 Sensor Matrix Relevance"
    mp1.font.size = Pt(12)
    mp1.font.bold = True
    mp1.font.color.rgb = COLOR_CYAN

    mp2 = mtf.add_paragraph()
    mp2.text = "Categorizing isolation origins directly guides sample preparation (e.g. blood lysate filtering vs. environmental water concentration) and prevents false-positive signals during field testing."
    mp2.font.size = Pt(12)
    mp2.font.color.rgb = COLOR_TEXT
    mp2.space_before = Pt(4)

    # Speaker Notes Slide 4
    slide4.notes_slide.notes_text_frame.text = (
        "Turning to our domain findings on Slide 4: we mapped isolation origins across three main ecological groups: "
        "clinical human specimens like blood and sputum for pathogens like Staphylococcus aureus; veterinary host samples "
        "like Rhodococcus equi from foal lung abscesses; and environmental wastewater samples like Corynebacterium glutamicum. "
        "Categorizing these origins is crucial for sensor data analysis because it helps us model background sample matrix noise—"
        "such as blood versus environmental water—preventing false-positive signals during field deployment."
    )

    # ==========================================
    # SLIDE 5: Growth Dynamics, Biosafety & SDS API Links
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Slide 5 / 5 • Safety & Culture Parameters", "Growth Dynamics, Biosafety & Safety Data Sheets", "Standardized culture growth requirements, biosafety containment levels, and direct REST API Safety Data Sheet integration.")

    # Top 2 Feature Cards
    top_cards = [
        ("🌡️ Growth Parameters & Profiles", "• Incubation Temp: 37°C clinical pathogen standard vs. 30°C environmental.\n• Culture Media: Trypticase Soy (#18/#260: 65%) and Brain Heart Infusion (#44: 20%).\n• Atmosphere: Aerobic (85%) vs. CO₂ / Microaerophilic (15%).", COLOR_AMBER),
        ("⚠️ Biosafety Rating & SDS API Links", "• BSL-1 Rating: Low-risk control strains (E. avium, C. glutamicum) for open bench.\n• BSL-2 Rating: Clinical pathogens (S. aureus, R. equi) requiring bio-hood protocols.\n• Direct SDS API: Dynamic REST links (/api/product/sds?atcc_number=...) for instant hazard access.", COLOR_CYAN)
    ]
    for idx, (gtitle, gdesc, gcolor) in enumerate(top_cards):
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 5.95), Inches(2.2), Inches(5.75), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = gcolor

        gtf = card.text_frame
        gtf.word_wrap = True
        gp1 = gtf.paragraphs[0]
        gp1.text = gtitle
        gp1.font.size = Pt(13)
        gp1.font.bold = True
        gp1.font.color.rgb = gcolor

        gp2 = gtf.add_paragraph()
        gp2.text = gdesc
        gp2.font.size = Pt(10)
        gp2.font.color.rgb = COLOR_MUTED
        gp2.space_before = Pt(6)

    # Data Table Excerpt (Bottom)
    table_shape = slide5.shapes.add_table(5, 7, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4))
    table = table_shape.table

    headers = ["ATCC #", "Organism Binomial", "Isolation Source", "Culture Medium", "Temp", "Atmosphere", "BSL"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(26, 37, 62)
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_CYAN

    table_data = [
        ["ATCC 14025", "Enterococcus avium", "Quality Control Record", "TSA w/ Sheep Blood (#260)", "37°C", "Aerobic", "BSL 1"],
        ["ATCC 8750", "Alcaligenes faecalis", "Preceptrol Culture", "Trypticase Soy (#18)", "37°C", "Aerobic", "BSL 1"],
        ["ATCC 6939", "Rhodococcus equi", "Foal Lung Abscess", "Brain Heart Infusion (#44)", "37°C", "Aerobic", "BSL 2"],
        ["ATCC 13032", "Corynebacterium glutamicum", "Sewage / Wastewater", "Nutrient Agar (#3)", "37°C", "Aerobic", "BSL 1"]
    ]
    for row_idx, row_data in enumerate(table_data):
        for col_idx, text_val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.text = text_val
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_TEXT

    # Speaker Notes Slide 5
    slide5.notes_slide.notes_text_frame.text = (
        "Finally, on Slide 5 covering growth parameters and biosafety: we extracted incubation temperatures, media formulations, "
        "and gaseous atmospheres across all 100 strains. Clinical pathogens cluster at 37°C while environmental strains grow "
        "at 30°C, primarily using Trypticase Soy Agar or Brain Heart Infusion broth. Every strain is also categorized by Biosafety "
        "Level—BSL-1 versus BSL-2—and paired with direct REST API links for official Safety Data Sheets for instant lab safety access. Thank you."
    )

    output_path = "/Users/madelynzhu/Desktop/bookish-giggle/100_Bug_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
